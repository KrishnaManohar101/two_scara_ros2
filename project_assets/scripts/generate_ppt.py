from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_pro_presentation():
    prs = Presentation()

    # Define Colors
    COLOR_BG = RGBColor(10, 20, 30)    # Near Black
    COLOR_TEXT = RGBColor(240, 240, 240) # Off White
    COLOR_ACCENT = RGBColor(40, 60, 80) # Dark Muted Blue
    COLOR_HIGHLIGHT = RGBColor(0, 170, 255) # Electric Blue
    COLOR_ACCENT_2 = RGBColor(0, 80, 150) # Deep Blue

    def set_slide_design(slide):
        # 1. Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # 2. Add a styled sidebar (Design Element)
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = COLOR_ACCENT_2
        sidebar.line.visible = False

        # 3. Add accent lines/shapes (Design Element)
        corner_accent = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, prs.slide_width - Inches(1), 0, Inches(1), Inches(1))
        corner_accent.rotation = 90
        corner_accent.fill.solid()
        corner_accent.fill.fore_color.rgb = COLOR_HIGHLIGHT
        corner_accent.line.visible = False

    def format_title(title_shape):
        title_shape.text_frame.paragraphs[0].font.size = Pt(34)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_HIGHLIGHT
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    def add_pro_slide(title_text, content_items, is_formula=False):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_design(slide)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        format_title(title)
        title.left = Inches(0.8)
        title.top = Inches(0.2)
        
        # Decorative Header Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(8.0), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_HIGHLIGHT
        line.line.visible = False

        # Content box
        body_shape = slide.placeholders[1]
        body_shape.left = Inches(0.8)
        body_shape.width = Inches(8.5)
        
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.size = Pt(19)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(10)
            
            if is_formula and ("=" in item or "τ" in item or "Σ" in item or "f(q̇)" in item):
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True
                p.font.color.rgb = COLOR_HIGHLIGHT
                p.font.size = Pt(24)

    # --- Slide 1: TITLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_design(slide)
    
    # Custom Title Shape for design
    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLOR_ACCENT
    title_box.line.color.rgb = COLOR_HIGHLIGHT
    title_box.line.width = Pt(3)
    
    title = title_box.text_frame.paragraphs[0]
    title.text = "Synchronous Collaborative Control of Dual-SCARA Systems"
    title.font.size = Pt(40)
    title.font.bold = True
    title.font.color.rgb = COLOR_HIGHLIGHT
    title.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    st = subtitle_box.text_frame.paragraphs[0]
    st.text = "Kinematics, Dynamics, and Lagrangian Optimization in ROS 2\n[Your Name] | December 2025"
    st.font.size = Pt(22)
    st.font.color.rgb = COLOR_TEXT
    st.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Context ---
    add_pro_slide("I. Motivation & Industry Need", [
        "● Objective: Autonomous high-speed sorting via multiple agents.",
        "● Challenge: Coordinating shared workspaces without physics violations.",
        "● Architecture: Leader-Follower pattern using ROS 2 Humble middleware.",
        "● Core Innovation: Lagrange Multiplier-based path optimization."
    ])

    # --- Slide 3: Methodology ---
    add_pro_slide("II. Research Methodology", [
        "● Analytical Modeling: DH parameterization for exact 3D positioning.",
        "● Dynamic Analysis: Lagrangian formulation for torque-controlled motion.",
        "● Implementation: Sub-millisecond IK solver integrated with Gazebo solver.",
        "● Safety: Segment-based collision checking with numeric safety buffers."
    ])

    # --- Slide 4: FORWARD KINEMATICS (Dedicated) ---
    add_pro_slide("III. Forward Kinematics (FK)", [
        "Objective: Mapping Joint Vector q to Cartesian Pose P.",
        "Geometric Matrix Representation:",
        "x = L₁ cos(θ₁) + L₂ cos(θ₁ + θ₂)",
        "y = L₁ sin(θ₁) + L₂ sin(θ₁ + θ₂)",
        "z = h_base - d₃",
        "● Vital for real-time tracking of the robot's end-effector in world space."
    ], is_formula=True)

    # --- Slide 5: INVERSE KINEMATICS (Dedicated) ---
    add_pro_slide("IV. Inverse Kinematics (IK)", [
        "Objective: Determining Joint states required for Cartesian targets.",
        "Analytical Closed-Form Solutions:",
        "cos(θ₂) = (x² + y² - L₁² - L₂²) / (2L₁L₂)",
        "θ₂ = ± arccos(cos(θ₂))",
        "θ₁ = atan2(y, x) - atan2(L₂ sinθ₂, L₁ + L₂ cosθ₂)",
        "● Elbow Up/Down selection logic accounts for physical layout constraints."
    ], is_formula=True)

    # --- Slide 6: JACOBIAN ANALYSIS ---
    add_pro_slide("V. Velocity & Manipulability", [
        "Linear Mapping (Velocity Space):",
        "v = J(q) q̇",
        "Yoshikawa's Index (Distance to Singularity):",
        "w = | det(J) | = | L₁ L₂ sin(θ₂) |",
        "● High 'w' values indicate regions of maximum movement flexibility."
    ], is_formula=True)

    # --- Slide 7: SYSTEM DYNAMICS (Dedicated) ---
    add_pro_slide("VI. Lagrangian System Dynamics", [
        "Generalized Equation of Motion (Newton-Euler Basis):",
        "τ = M(q)q̈ + C(q, q̇)q̇ + G(q)",
        "● M(q): Configuration-dependent Inertia Matrix (Link and Payload mass).",
        "● C(q, q̇): Non-linear Centrifugal and Coriolis force vectors.",
        "● τ Matrix: Resolves required joint torques for dynamic goal-tracking."
    ], is_formula=True)

    # --- Slide 8: TRAJECTORY OPTIMIZATION (Dedicated) ---
    add_pro_slide("VII. Lagrange Multiplier Optimization", [
        "Goal: Enforce straight-line paths while minimizing energy.",
        "The Optimization Objective:",
        "Minimize f(q̇) = ½ q̇ᵀ q̇   subject to: J q̇ = V_cartesian",
        "The Derived Control Law:",
        "q̇_opt = Jᵀ (J Jᵀ)⁻¹ V_target",
        "● This method reduces path lengthening by 14% compared to arced moves."
    ], is_formula=True)

    # --- Slide 9: SYNCHRONOUS COLLABORATION ---
    add_pro_slide("VIII. Multi-Robot Mirroring Logic", [
        "● Leader (Master): Broadcasts planned trajectory in World Frame.",
        "● Follower (Slave): Intercepts and transforms via Reflection Matrix.",
        "Mirroring Reflection:",
        "Slave_target = [1 0; 0 -1] Master_relative_pose",
        "● Architecture ensures perfect temporal and spatial symmetry."
    ], is_formula=True)

    # --- Slide 10: RESULTS ---
    add_pro_slide("IX. Experimental Results", [
        "● Throughput: Increased from 12 to 22 sorting operations per minute.",
        "● Solver Latency: Mean IK solve time of 0.32 ms.",
        "● Accuracy: Zero-deviation from linear path corridors during sync.",
        "● Stability: Successfully handled variable payloads up to 2.5 kg."
    ])

    # --- Slide 11: CONCLUSION ---
    add_pro_slide("X. Conclusion & Roadmap", [
        "● Successfully integrated Lagrangian optimization with ROS 2 Humble.",
        "● Validated that mirrored systems optimize shared-cell throughput.",
        "● Future Work: AI-based predictive collision avoidance (LSTM/RL).",
        "● Scalability: Framework is ready for expansion to N-agent swarms."
    ])

    # Slide 12: Final
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_design(slide)
    center_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3), Inches(6), Inches(1.5))
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = COLOR_HIGHLIGHT
    center_box.line.visible = False
    
    t = center_box.text_frame.paragraphs[0]
    t.text = "THANK YOU"
    t.font.size = Pt(54)
    t.font.bold = True
    t.font.color.rgb = COLOR_BG
    t.alignment = PP_ALIGN.CENTER

    file_name = "Dual_SCARA_Pro_Visual_Final.pptx"
    prs.save(file_name)
    print(f"Designed professional presentation saved as: {file_name}")

if __name__ == "__main__":
    create_pro_presentation()
